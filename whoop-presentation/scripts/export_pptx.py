#!/usr/bin/env python3
"""
export_pptx.py — WHOOP HTML → PPTX exporter

Usage:
    python scripts/export_pptx.py --html final-slides.html --out final-slides.pptx

Parses each reveal.js <section>, detects slide type from CSS classes,
and builds a WHOOP-branded PPTX. Never crashes — bad slides get a fallback.
"""

import argparse
import re
import sys
import warnings
from pathlib import Path

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Brand constants
def px(n): return Emu(n * 9525)

SLIDE_W  = px(1280)
SLIDE_H  = px(720)
MARGIN_X = px(48)
BODY_TOP = px(92)
BODY_W   = px(1184)
COL_W    = px(570)
COL_GAP  = px(44)
HEADER_H = px(52)

C_BG_DARK = RGBColor(0x10, 0x15, 0x18)
C_BG_DIV  = RGBColor(0x0D, 0x12, 0x15)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED   = RGBColor(0xAA, 0xAA, 0xAA)
C_LABEL   = RGBColor(0x66, 0x66, 0x66)
C_TEAL    = RGBColor(0x00, 0xF1, 0x9F)
C_BLUE    = RGBColor(0x00, 0x93, 0xE7)
C_PINK    = RGBColor(0xE8, 0x79, 0xA0)
C_DIVIDER = RGBColor(0x22, 0x2A, 0x2E)


# Low-level shape helpers
def _rect(slide, left, top, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def _textbox(slide, left, top, w, h, text, size,
             bold=False, color=None, caps=False, align=PP_ALIGN.LEFT, proxima=False):
    if not text:
        return None
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Proxima Nova" if proxima else "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    if caps:
        run.font.all_caps = True
    return tb


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _add_header(slide, deck_title, year, proxima=False):
    _rect(slide, px(0), px(0), SLIDE_W, HEADER_H, C_BG_DARK)
    _rect(slide, MARGIN_X, HEADER_H, BODY_W, px(1), C_DIVIDER)
    _textbox(slide, MARGIN_X, px(18), px(80), px(20),
             "WHOOP", 11, bold=True, color=C_WHITE, caps=True, proxima=proxima)
    _textbox(slide, px(0), px(18), SLIDE_W - MARGIN_X, px(20),
             f"WHOOP  ·  {year}  ·  {deck_title.upper()}  ·  CONFIDENTIAL",
             9, color=C_LABEL, caps=True, align=PP_ALIGN.RIGHT, proxima=proxima)


# Shared builder helpers
def _new_slide(prs, meta, bg=None):
    """Add blank slide, set background, add header. Returns (slide, proxima)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, bg or C_BG_DARK)
    proxima = meta.get('proxima_nova_available', False)
    _add_header(slide, meta['title'], meta['year'], proxima=proxima)
    return slide, proxima


def _render_preamble(slide, data, y, proxima, *, with_subtitle=False):
    """Render section label + title (+ subtitle). Returns updated y."""
    sec = data.get('section_label', '')
    if sec:
        _textbox(slide, MARGIN_X, y, BODY_W, px(18),
                 sec.upper(), 10, bold=True, color=C_TEAL, caps=True, proxima=proxima)
        y += px(20)
    title = data.get('title', '')
    if title:
        _textbox(slide, MARGIN_X, y, BODY_W, px(42),
                 title.upper(), 26, bold=True, color=C_WHITE, caps=True, proxima=proxima)
        y += px(44) if with_subtitle else px(52)
    if with_subtitle:
        sub = data.get('subtitle', '')
        if sub:
            _textbox(slide, MARGIN_X, y, BODY_W, px(22), sub, 13, color=C_MUTED, proxima=proxima)
            y += px(28)
    return y


def _render_col_bullets(slide, x, y, bullets, w, proxima):
    """Render a list of bullet dicts into a column. Returns updated y."""
    for b in bullets:
        text = b.get('text', str(b)) if isinstance(b, dict) else str(b)
        _textbox(slide, x, y, w, px(28), '—  ' + text, 13, bold=True, proxima=proxima)
        y += px(34)
    return y


def _render_link_items(slide, items, y, proxima):
    """Render label / description / url triples. Returns updated y."""
    for item in items:
        label, desc, url = item.get('label',''), item.get('description',''), item.get('url','')
        if label:
            _textbox(slide, MARGIN_X, y, BODY_W, px(18), label, 12, bold=True, proxima=proxima)
            y += px(20)
        if desc and desc != label:
            _textbox(slide, MARGIN_X, y, BODY_W, px(18), desc, 12, color=C_MUTED, proxima=proxima)
            y += px(20)
        if url:
            _textbox(slide, MARGIN_X, y, BODY_W, px(18), url, 11, color=C_BLUE, proxima=proxima)
            y += px(28)
    return y


# HTML parsing helpers
def _text(tag):
    return re.sub(r'\s+', ' ', tag.get_text(separator=' ')).strip() if tag else ''

def _has_class(tag, *names):
    return any(c in (tag.get('class') or []) for c in names)

def _find_class(tag, *names):
    for name in names:
        found = tag.find(class_=name)
        if found:
            return found
    return None

def _parse_bullets(ul):
    if ul is None:
        return []
    items = []
    for li in ul.find_all('li', recursive=False):
        text = re.sub(r'^[—·\-]\s*', '', _text(li)).strip()
        if text:
            items.append({'text': text, 'sub': _has_class(li, 'sub')})
    return items

def _parse_callouts(container):
    if container is None:
        return []
    return [{'style': 'warning' if _has_class(t, 'callout-warn') else 'info', 'text': _text(t)}
            for t in container.find_all(class_=['callout', 'callout-warn'])]


# Slide type detection
def _detect_type(slide_div):
    if slide_div is None:
        return 'content'
    classes = slide_div.get('class') or []
    if 'slide-title'   in classes: return 'title'
    if 'slide-divider' in classes: return 'divider'
    body = slide_div.find(class_='slide-body') or slide_div
    if body.find(class_=['challenge-col', 'col-label-pain', 'col-label-ask']): return 'challenge'
    if body.find(class_='two-col'): return 'two_col'
    if len(body.find_all('a')) >= 2 and not body.find(class_='b'): return 'reference'
    return 'content'


# Per-type HTML parsers
def _parse_title(div, fallback):
    return {'type': 'title',
            'title':    _text(div.find('h1')) or fallback,
            'subtitle': _text(div.find(class_='subtitle') or div.find('p'))}

def _parse_divider(div):
    body = div.find(class_='slide-body') or div
    return {'type': 'divider', 'section': _text(body.find('h2'))}

def _parse_content(div):
    body = div.find(class_='slide-body') or div
    bullets = _parse_bullets(body.find('ul', class_='b'))

    flow = body.find(class_='flow')
    if flow:
        for step in flow.find_all(class_='flow-step'):
            txt = _text(step.find(class_='flow-text'))
            num = _text(step.find(class_='flow-num'))
            if txt:
                bullets.append({'text': f"{num}. {txt}".strip('. '), 'sub': False})

    for tbl in body.find_all('table', class_='t'):
        hdrs = [_text(th) for th in tbl.find_all('th')]
        if hdrs:
            bullets.append({'text': ' | '.join(hdrs), 'sub': False})
        for row in tbl.find_all('tr'):
            cells = [_text(td) for td in row.find_all('td')]
            if cells:
                bullets.append({'text': ' | '.join(cells), 'sub': True})

    for stat in body.find_all(class_='stat-num'):
        label = stat.find_next_sibling(class_='stat-label')
        t = _text(stat) + ('  ' + _text(label) if label else '')
        if t:
            bullets.insert(0, {'text': t, 'sub': False})

    return {
        'type': 'content',
        'section_label': _text(_find_class(body, 'section-label')),
        'title':    _text(body.find(class_='content-title') or body.find('h2')),
        'subtitle': _text(_find_class(body, 'content-subtitle')),
        'bullets':  bullets,
        'callouts': _parse_callouts(body),
        'links':    [{'label': _text(a), 'url': a.get('href',''), 'description': ''}
                     for a in body.find_all('a')],
    }

def _parse_two_col_base(div, col_key, col_parser):
    body = div.find(class_='slide-body') or div
    two_col = body.find(class_='two-col')
    cols = [c for c in two_col.children if isinstance(c, Tag)] if two_col else []
    if not cols:
        cols = body.find_all(class_='challenge-col')
    return {
        'section_label': _text(_find_class(body, 'section-label')),
        'title': _text(body.find(class_='content-title') or body.find('h2')),
        'left':  col_parser(cols[0] if len(cols) > 0 else None),
        'right': col_parser(cols[1] if len(cols) > 1 else None),
    }

def _parse_challenge(div):
    def _col(tag):
        if tag is None: return {'label': '', 'bullets': []}
        return {'label':   _text(_find_class(tag, 'col-label', 'col-label-pain', 'col-label-ask')),
                'bullets': _parse_bullets(tag.find('ul', class_='b'))}
    return {'type': 'challenge', **_parse_two_col_base(div, 'label', _col)}

def _parse_two_col(div):
    def _col(tag):
        if tag is None: return {'heading': '', 'bullets': []}
        return {'heading': _text(tag.find(['h3','h4']) or tag.find('strong')),
                'bullets': _parse_bullets(tag.find('ul', class_='b'))}
    return {'type': 'two_col', **_parse_two_col_base(div, 'heading', _col)}

def _parse_reference(div):
    body = div.find(class_='slide-body') or div
    return {
        'type': 'reference',
        'title': _text(body.find(class_='content-title') or body.find('h2')),
        'items': [{'label': _text(a), 'url': a.get('href',''),
                   'description': _text(a.find_previous_sibling(['p','span','div']))}
                  for a in body.find_all('a')],
    }


# HTML → slide data
def parse_html(html_path, verbose=False):
    html = html_path.read_text(encoding='utf-8', errors='replace')
    soup = BeautifulSoup(html, 'html.parser')

    title_tag = soup.find('title')
    deck_title = _text(title_tag) if title_tag else html_path.stem.replace('-', ' ').title()
    year_match = re.search(r'WHOOP\s*[·•]\s*(\d{4})', html, re.IGNORECASE)
    meta = {'title': deck_title, 'year': year_match.group(1) if year_match else '2026',
            'proxima_nova_available': False}

    sections = soup.find_all('section')
    if not sections:
        warnings.warn("No <section> tags found — is this a reveal.js file?")
        return meta, []

    parsers = {'title': _parse_title, 'divider': _parse_divider, 'challenge': _parse_challenge,
               'two_col': _parse_two_col, 'reference': _parse_reference}
    slides = []

    for i, section in enumerate(sections, 1):
        try:
            div = section.find(class_='slide')
            if not div:
                if verbose: print(f"  ⚠️  Slide {i}: no .slide div — skipping")
                continue
            stype = _detect_type(div)
            if stype == 'title':
                data = _parse_title(div, deck_title)
            else:
                data = parsers.get(stype, _parse_content)(div)
            slides.append(data)
            if verbose:
                print(f"  Slide {i}: [{stype}]  {data.get('title') or data.get('section','')}")
        except Exception as exc:
            warnings.warn(f"Slide {i} parse error: {exc} — using raw text fallback")
            slides.append({'type': 'content', 'title': f'Slide {i}',
                           'subtitle': '(parse error)', 'bullets': [], 'callouts': [],
                           'links': [], 'notes': _text(section)[:400]})
    return meta, slides


# PPTX slide builders
def _build_title(prs, data, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_DARK)
    proxima = meta.get('proxima_nova_available', False)
    _rect(slide, px(0), px(0), px(4), SLIDE_H, C_TEAL)
    _textbox(slide, px(64), px(200), px(1100), px(200),
             data.get('title','').upper(), 44, bold=True, caps=True, proxima=proxima)
    _rect(slide, px(64), px(372), px(120), px(3), C_TEAL)
    _textbox(slide, px(64), px(390), px(800), px(40),
             data.get('subtitle', meta.get('year','')), 16, color=C_MUTED, proxima=proxima)
    _textbox(slide, MARGIN_X, px(680), BODY_W, px(24),
             'CONFIDENTIAL', 9, bold=True, color=C_LABEL, caps=True,
             align=PP_ALIGN.RIGHT, proxima=proxima)

def _build_divider(prs, data, meta):
    slide, proxima = _new_slide(prs, meta, bg=C_BG_DIV)
    _rect(slide, MARGIN_X, px(260), px(3), px(80), C_TEAL)
    _textbox(slide, px(64), px(260), BODY_W, px(180),
             data.get('section','').upper(), 42, bold=True, caps=True, proxima=proxima)

def _build_content(prs, data, meta):
    slide, proxima = _new_slide(prs, meta)
    y = _render_preamble(slide, data, BODY_TOP, proxima, with_subtitle=True)

    for b in data.get('bullets', []):
        text  = b.get('text', str(b)) if isinstance(b, dict) else str(b)
        is_sub = b.get('sub', False) if isinstance(b, dict) else False
        _textbox(slide, MARGIN_X + (px(24) if is_sub else 0), y,
                 BODY_W - (px(24) if is_sub else 0), px(28),
                 ('·  ' if is_sub else '—  ') + text,
                 12 if is_sub else 14, bold=not is_sub,
                 color=C_MUTED if is_sub else C_WHITE, proxima=proxima)
        y += px(24 if is_sub else 30)

    for c in data.get('callouts', []):
        _rect(slide, MARGIN_X, y, px(3), px(40),
              C_BLUE if c.get('style') == 'info' else C_PINK)
        _textbox(slide, MARGIN_X + px(12), y + px(6), BODY_W - px(16), px(30),
                 c.get('text',''), 13, proxima=proxima)
        y += px(52)

    _render_link_items(slide, data.get('links', []), y, proxima)

    if data.get('notes'):
        slide.notes_slide.notes_text_frame.text = data['notes']

def _build_challenge(prs, data, meta):
    slide, proxima = _new_slide(prs, meta)
    y = _render_preamble(slide, data, BODY_TOP, proxima)
    lx, rx = MARGIN_X, MARGIN_X + COL_W + COL_GAP
    left, right = data.get('left', {}), data.get('right', {})

    _rect(slide, lx, y, px(180), px(22), RGBColor(0x2A, 0x14, 0x1D))
    _textbox(slide, lx + px(8), y + px(3), px(164), px(18),
             left.get('label', 'Current State').upper(), 9, bold=True, color=C_PINK, caps=True, proxima=proxima)
    _rect(slide, rx, y, px(180), px(22), RGBColor(0x0A, 0x1A, 0x2A))
    _textbox(slide, rx + px(8), y + px(3), px(164), px(18),
             right.get('label', "What We're Looking For").upper(), 9, bold=True, color=C_BLUE, caps=True, proxima=proxima)
    y += px(36)
    _rect(slide, MARGIN_X + COL_W + COL_GAP // 2, y, px(1), px(400), C_DIVIDER)
    _render_col_bullets(slide, lx, y, left.get('bullets', []),  COL_W, proxima)
    _render_col_bullets(slide, rx, y, right.get('bullets', []), COL_W, proxima)

def _build_two_col(prs, data, meta):
    slide, proxima = _new_slide(prs, meta)
    y = _render_preamble(slide, data, BODY_TOP, proxima)
    lx, rx = MARGIN_X, MARGIN_X + COL_W + COL_GAP
    left, right = data.get('left', {}), data.get('right', {})

    if left.get('heading'):
        _textbox(slide, lx, y, COL_W, px(22), left['heading'].upper(),
                 11, bold=True, color=C_TEAL, caps=True, proxima=proxima)
    if right.get('heading'):
        _textbox(slide, rx, y, COL_W, px(22), right['heading'].upper(),
                 11, bold=True, color=C_TEAL, caps=True, proxima=proxima)
    y += px(30)
    _render_col_bullets(slide, lx, y, left.get('bullets', []),  COL_W, proxima)
    _render_col_bullets(slide, rx, y, right.get('bullets', []), COL_W, proxima)

def _build_reference(prs, data, meta):
    slide, proxima = _new_slide(prs, meta)
    y = BODY_TOP
    _textbox(slide, MARGIN_X, y, BODY_W, px(42),
             data.get('title', 'REFERENCES').upper(), 26, bold=True, caps=True, proxima=proxima)
    _render_link_items(slide, data.get('items', []), y + px(52), proxima)


BUILDERS = {
    'title':     _build_title,
    'divider':   _build_divider,
    'content':   _build_content,
    'challenge': _build_challenge,
    'two_col':   _build_two_col,
    'reference': _build_reference,
}


# Entry point
def main():
    parser = argparse.ArgumentParser(description="Export reveal.js HTML → PPTX (WHOOP branding)")
    parser.add_argument('--html', default='final-slides.html')
    parser.add_argument('--out',  default='final-slides.pptx')
    parser.add_argument('--verbose', action='store_true')
    args = parser.parse_args()

    html_path = Path(args.html)
    if not html_path.exists():
        print(f"❌  Not found: {html_path}")
        sys.exit(1)

    print(f"📖  Parsing: {html_path}")
    meta, slides = parse_html(html_path, verbose=args.verbose)

    if not slides:
        print("❌  No slides parsed. Check that the file uses the WHOOP reveal.js structure.")
        sys.exit(1)

    print(f"📊  Building PPTX: {len(slides)} slides — \"{meta['title']}\"")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(slides, 1):
        stype   = slide_data.get('type', 'content')
        builder = BUILDERS.get(stype, _build_content)
        try:
            builder(prs, slide_data, meta)
        except Exception as exc:
            warnings.warn(f"Slide {i} build error ({stype}): {exc} — inserting blank")
            blank = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(blank, C_BG_DARK)
            _add_header(blank, meta['title'], meta['year'])
            _textbox(blank, MARGIN_X, BODY_TOP, BODY_W, px(40),
                     f"Slide {i} — render error: {exc}", 14, color=C_PINK)

    out_path = Path(args.out)
    prs.save(str(out_path))
    print(f"✅  Saved: {out_path}  ({out_path.stat().st_size // 1024} KB, {len(slides)} slides)")
    if not meta.get('proxima_nova_available'):
        print("ℹ️   Font: Helvetica Neue (Proxima Nova not detected)")


if __name__ == '__main__':
    main()
